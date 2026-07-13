from __future__ import annotations

from html import escape
from pathlib import Path


def build_report_file(root: Path, folder_id: str, report_id: str, file_format: str, title: str, body: str) -> Path:
    folder_dir = root / _safe_name(folder_id)
    folder_dir.mkdir(parents=True, exist_ok=True)
    fmt = file_format.lower()
    path = folder_dir / f"{_safe_name(report_id)}.{fmt}"

    if fmt == "pdf":
        path.write_bytes(_minimal_pdf(title, body))
    elif fmt == "html":
        _write_html(path, title, body)
    elif fmt == "docx":
        _write_docx(path, title, body)
    elif fmt == "pptx":
        _write_pptx(path, title, body)
    elif fmt == "xlsx":
        _write_xlsx(path, title, body)
    else:
        raise ValueError(f"Unsupported report format: {file_format}")
    return path


def _write_html(path: Path, title: str, body: str) -> None:
    paragraphs = "\n".join(
        f"<p>{escape(paragraph.strip())}</p>"
        for paragraph in body.split("\n\n")
        if paragraph.strip()
    )
    document = f"""<!doctype html>
<html lang="en">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <title>{escape(title)}</title>
  <style>
    body {{ max-width: 820px; margin: 48px auto; padding: 0 24px; color: #171717; background: #fff; font: 16px/1.6 Arial, sans-serif; }}
    h1 {{ font-size: 30px; line-height: 1.2; margin-bottom: 32px; }}
    p {{ margin: 0 0 18px; white-space: pre-wrap; }}
  </style>
</head>
<body>
  <h1>{escape(title)}</h1>
  {paragraphs}
</body>
</html>
"""
    path.write_text(document, encoding="utf-8")

def _write_docx(path: Path, title: str, body: str) -> None:
    try:
        from docx import Document

        doc = Document()
        doc.add_heading(title, level=1)
        for paragraph in body.split("\n\n"):
            doc.add_paragraph(paragraph.strip())
        doc.save(path)
    except Exception:
        path.write_text(f"{title}\n\n{body}", encoding="utf-8")


def _write_pptx(path: Path, title: str, body: str) -> None:
    try:
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        slide.placeholders[1].text = _clip(body, 600)

        for chunk in _chunks(body, 700)[:5]:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Analysis"
            slide.placeholders[1].text = chunk
        prs.save(path)
    except Exception:
        path.write_text(f"{title}\n\n{body}", encoding="utf-8")


def _write_xlsx(path: Path, title: str, body: str) -> None:
    try:
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws.title = "Report"
        ws["A1"] = title
        row = 3
        for paragraph in body.split("\n"):
            ws.cell(row=row, column=1).value = paragraph
            row += 1
        wb.save(path)
    except Exception:
        path.write_text(f"{title}\n\n{body}", encoding="utf-8")


def _minimal_pdf(title: str, body: str) -> bytes:
    text = f"{title}\n\n{body}"
    lines = []
    y = 760
    for raw in text.splitlines():
        for line in _wrap(raw, 88):
            lines.append(f"BT /F1 10 Tf 50 {y} Td ({_pdf_escape(line)}) Tj ET")
            y -= 14
            if y < 50:
                break
        if y < 50:
            break

    stream = "\n".join(lines).encode("latin-1", errors="replace")
    objects = [
        b"1 0 obj << /Type /Catalog /Pages 2 0 R >> endobj\n",
        b"2 0 obj << /Type /Pages /Kids [3 0 R] /Count 1 >> endobj\n",
        b"3 0 obj << /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >> endobj\n",
        b"4 0 obj << /Type /Font /Subtype /Type1 /BaseFont /Helvetica >> endobj\n",
        f"5 0 obj << /Length {len(stream)} >> stream\n".encode("ascii") + stream + b"\nendstream endobj\n",
    ]
    pdf = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for obj in objects:
        offsets.append(len(pdf))
        pdf.extend(obj)
    xref = len(pdf)
    pdf.extend(f"xref\n0 {len(objects)+1}\n0000000000 65535 f \n".encode("ascii"))
    for offset in offsets[1:]:
        pdf.extend(f"{offset:010d} 00000 n \n".encode("ascii"))
    pdf.extend(f"trailer << /Size {len(objects)+1} /Root 1 0 R >>\nstartxref\n{xref}\n%%EOF".encode("ascii"))
    return bytes(pdf)


def _wrap(text: str, width: int) -> list[str]:
    if not text:
        return [""]
    words = text.split()
    lines = []
    current = ""
    for word in words:
        if len(current) + len(word) + 1 > width:
            lines.append(current)
            current = word
        else:
            current = f"{current} {word}".strip()
    if current:
        lines.append(current)
    return lines or [""]


def _chunks(text: str, size: int) -> list[str]:
    return [text[i : i + size] for i in range(0, len(text), size)] or [text]


def _clip(text: str, size: int) -> str:
    return text if len(text) <= size else text[: size - 3] + "..."


def _pdf_escape(text: str) -> str:
    return text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")


def _safe_name(value: str) -> str:
    return "".join(ch for ch in value if ch.isalnum() or ch in ("-", "_")) or "artifact"

